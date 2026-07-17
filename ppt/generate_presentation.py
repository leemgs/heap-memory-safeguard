#!/usr/bin/env python3
"""Generate the Korean HMS TECS presentation deck."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ppt" / "HMS_TECS_presentation.pptx"

NAVY = RGBColor(15, 27, 45)
NAVY2 = RGBColor(26, 43, 67)
TEAL = RGBColor(29, 183, 168)
CYAN = RGBColor(70, 202, 220)
ORANGE = RGBColor(246, 166, 73)
RED = RGBColor(230, 91, 96)
WHITE = RGBColor(247, 249, 252)
MUTED = RGBColor(169, 184, 204)
GRID = RGBColor(64, 82, 108)
FONT = "Noto Sans CJK KR"


def add_text(slide, text, x, y, w, h, size=24, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=NAVY2, line=GRID, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1)
    return s


def add_bg(slide, num):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    add_rect(slide, 0, 0, 13.333, 0.08, TEAL, TEAL)
    add_text(slide, f"{num:02d}", 12.55, 7.05, 0.45, 0.25, 10, MUTED,
             align=PP_ALIGN.RIGHT)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.65, 0.34, 12.0, 0.55, 28, WHITE, True)
    if subtitle:
        add_text(slide, subtitle, 0.68, 0.94, 11.8, 0.35, 12, MUTED)


def add_bullets(slide, items, x, y, w, h, size=20, color=WHITE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
        p.text = "• " + p.text
    return box


def card(slide, title, body, x, y, w, h, accent=TEAL):
    add_rect(slide, x, y, w, h, NAVY2, GRID, True)
    add_rect(slide, x, y, 0.08, h, accent, accent)
    add_text(slide, title, x + 0.3, y + 0.22, w - 0.5, 0.4, 18, accent, True)
    add_text(slide, body, x + 0.3, y + 0.72, w - 0.55, h - 0.9, 14, WHITE)


def connector(slide, x1, y1, x2, y2, color=MUTED):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    c.line.color.rgb = color
    c.line.width = Pt(2)
    return c


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank); add_bg(s, 1)
    add_text(s, "Heap Memory Safeguard", 0.72, 1.25, 11.9, 0.9, 36, WHITE, True)
    add_text(s, "Android cross-layer memory coordination을 위한\nsimulation-backed design proposal",
             0.76, 2.25, 10.8, 1.25, 24, TEAL, True)
    add_rect(s, 0.78, 4.05, 5.3, 1.2, NAVY2, GRID, True)
    add_text(s, "Geunsik Lim  ·  Sungkyunkwan University",
             1.08, 4.34, 4.7, 0.35, 17, WHITE, True)
    add_text(s, "leemgs@g.skku.edu  ·  Corresponding author",
             1.08, 4.77, 4.7, 0.3, 12, MUTED)
    add_text(s, "ACM TECS manuscript presentation", 0.8, 6.5, 4.5, 0.3, 12, MUTED)

    # 2. Takeaway
    s = prs.slides.add_slide(blank); add_bg(s, 2)
    add_title(s, "오늘의 핵심 메시지", "주장 범위와 증거 범위를 일치시킨 개정본")
    card(s, "문제", "ART managed heap과 native heap, Linux page pressure가 서로 다른 관측 단위를 사용한다.",
         0.7, 1.45, 3.85, 2.0, ORANGE)
    card(s, "아이디어", "RSS occupancy와 online net-growth를 결합해 bounded memcg reclaim 요청을 만든다.",
         4.75, 1.45, 3.85, 2.0, TEAL)
    card(s, "현재 증거", "Deterministic simulation으로 controller trade-off를 검증한다. 실제 Android 성능은 주장하지 않는다.",
         8.8, 1.45, 3.85, 2.0, CYAN)
    add_rect(s, 0.7, 4.05, 11.95, 1.65, NAVY2, TEAL, True)
    add_text(s, "정직한 결론", 1.05, 4.34, 2.0, 0.4, 20, TEAL, True)
    add_text(s, "HMS는 배포 완료 시스템이 아니라, AOSP/kernel 구현과 device study를 위한 감사 가능한 설계 기반이다.",
             3.0, 4.25, 9.1, 0.75, 21, WHITE, True)

    # 3. Problem
    s = prs.slides.add_slide(blank); add_bg(s, 3)
    add_title(s, "문제: 세 개의 서로 다른 관측 창")
    card(s, "ART", "Managed object\nGC phase / pause\nNative ownership은 불완전", 0.8, 1.45, 3.35, 3.35, TEAL)
    card(s, "Native allocator", "malloc/free\nallocator cache\nmapping / reuse / lifetime", 4.98, 1.45, 3.35, 3.35, ORANGE)
    card(s, "Linux kernel", "RSS / memcg charge\nPSI stall\npage reclaim / refault", 9.15, 1.45, 3.35, 3.35, CYAN)
    connector(s, 4.15, 3.1, 4.98, 3.1)
    connector(s, 8.33, 3.1, 9.15, 3.1)
    add_text(s, "Visibility gap", 2.15, 5.25, 3.4, 0.4, 19, ORANGE, True, PP_ALIGN.CENTER)
    add_text(s, "Temporal gap", 7.75, 5.25, 3.4, 0.4, 19, CYAN, True, PP_ALIGN.CENTER)
    add_text(s, "※ native allocation ≠ page fault ≠ RSS growth", 3.6, 6.15, 6.2, 0.35, 17, MUTED, True, PP_ALIGN.CENTER)

    # 4. Existing mechanisms
    s = prs.slides.add_slide(blank); add_bg(s, 4)
    add_title(s, "기존 Android/Linux 메커니즘을 대체하지 않는다")
    mechanisms = [
        ("LMKD + PSI", "pressure와 importance 기반 kill"),
        ("memory.reclaim", "cgroup-scoped proactive reclaim"),
        ("DAMON/DAMOS", "cold-region monitoring + action"),
        ("MGLRU", "working-set / thrashing control"),
        ("CachedAppOptimizer / mmd", "compaction, zRAM writeback/prefetch"),
        ("heapprofd", "sampled allocation attribution"),
    ]
    for i, (t, b) in enumerate(mechanisms):
        x = 0.75 + (i % 3) * 4.15
        y = 1.45 + (i // 3) * 2.15
        card(s, t, b, x, y, 3.75, 1.65, [TEAL, CYAN, ORANGE][i % 3])
    add_text(s, "HMS의 질문: “기존 actuator를 언제, 누구에게, 얼마나 요청할 것인가?”",
             1.0, 6.05, 11.3, 0.45, 21, WHITE, True, PP_ALIGN.CENTER)

    # 5. Architecture
    s = prs.slides.add_slide(blank); add_bg(s, 5)
    add_title(s, "제안 아키텍처와 artifact boundary")
    card(s, "Observer model", "managed activity\nsampled native activity\nRSS / memcg / PSI",
         0.75, 1.5, 3.3, 2.6, TEAL)
    card(s, "HMS controller", "online normalization\nθ₁ / θ₂ hysteresis\nbudget + cooldown",
         5.0, 1.5, 3.3, 2.6, ORANGE)
    card(s, "Existing actuator", "memory.reclaim\nDAMON/MGLRU 후보\nrefault feedback",
         9.25, 1.5, 3.3, 2.6, CYAN)
    connector(s, 4.05, 2.8, 5.0, 2.8, TEAL)
    connector(s, 8.3, 2.8, 9.25, 2.8, ORANGE)
    add_rect(s, 0.75, 4.8, 11.8, 1.05, NAVY2, RED, True)
    add_text(s, "현재 repo: Python simulation", 1.1, 5.08, 3.2, 0.35, 18, RED, True)
    add_text(s, "미포함: system_server · ART hook · kernel patch · SELinux · device trace",
             4.25, 5.08, 7.8, 0.35, 17, WHITE)

    # 6. Formula
    s = prs.slides.add_slide(blank); add_bg(s, 6)
    add_title(s, "Controller state: heuristic score와 해석 한계")
    add_rect(s, 0.75, 1.45, 7.4, 1.5, NAVY2, TEAL, True)
    add_text(s, "Hᵢ(t) = Rᵢ(t)/Lᵢ(t) + α · ĝᵢ(t)/Gᵢ(t)",
             1.05, 1.84, 6.8, 0.5, 27, WHITE, True, PP_ALIGN.CENTER)
    card(s, "R / L", "현재 resident occupancy", 0.75, 3.35, 3.45, 1.55, TEAL)
    card(s, "ĝ", "smoothed net resident growth", 4.45, 3.35, 3.45, 1.55, ORANGE)
    card(s, "G", "trailing p95 online scale\n미래 sample 사용 없음", 8.15, 3.35, 3.45, 1.55, CYAN)
    add_rect(s, 8.55, 1.45, 3.7, 1.5, NAVY2, RED, True)
    add_text(s, "α는 prediction horizon이 아니라\n무차원 weight",
             8.85, 1.75, 3.1, 0.85, 18, RED, True, PP_ALIGN.CENTER)
    add_text(s, "향후 대안: T_pressure = headroom / max(net growth, ε)",
             1.4, 5.7, 10.4, 0.45, 20, MUTED, True, PP_ALIGN.CENTER)

    # 7. Control loop
    s = prs.slides.add_slide(blank); add_bg(s, 7)
    add_title(s, "Bounded control loop")
    steps = [
        ("1", "Observe", TEAL),
        ("2", "Normalize", CYAN),
        ("3", "θ₁ watch", ORANGE),
        ("4", "θ₂ request", RED),
        ("5", "Budgeted reclaim", TEAL),
        ("6", "Cooldown / feedback", CYAN),
    ]
    for i, (n, label, col) in enumerate(steps):
        x = 0.62 + i * 2.08
        add_rect(s, x, 2.0, 1.7, 1.55, NAVY2, col, True)
        add_text(s, n, x + 0.15, 2.18, 0.35, 0.35, 18, col, True)
        add_text(s, label, x + 0.2, 2.62, 1.3, 0.55, 15, WHITE, True, PP_ALIGN.CENTER)
        if i < 5:
            connector(s, x + 1.7, 2.78, x + 2.08, 2.78, MUTED)
    add_bullets(s, [
        "per-cgroup / global budget로 cycle당 work 상한",
        "GC-active 및 render-deadline window에서는 reclaim 지연",
        "foreground class 우선 + class 내부 deficit round-robin",
        "process exit / cgroup identity 재검증",
    ], 1.05, 4.25, 11.2, 1.9, 18)

    # 8. Artifact
    s = prs.slides.add_slide(blank); add_bg(s, 8)
    add_title(s, "재현 가능한 simulation artifact")
    card(s, "Workloads", "W1 camera-like bursts\nW2 game streaming\nW3 inference tensors\nW4 app switching",
         0.75, 1.4, 3.5, 3.5, TEAL)
    card(s, "Determinism", "seeded telemetry\nindependent noise streams\nworkload별 controller reset\n동일 trace 비교",
         4.9, 1.4, 3.5, 3.5, CYAN)
    card(s, "Outputs", "summary.csv\nparameter sweeps\nPNG figures\nconfiguration + seed",
         9.05, 1.4, 3.5, 3.5, ORANGE)
    add_rect(s, 1.35, 5.5, 10.6, 0.75, NAVY2, TEAL, True)
    add_text(s, "python scripts/run_experiments.py --seed 7 --outdir results",
             1.65, 5.72, 10.0, 0.3, 16, WHITE, True, PP_ALIGN.CENTER)

    # 9. Results
    s = prs.slides.add_slide(blank); add_bg(s, 9)
    add_title(s, "Simulation 결과: ‘모두 개선’이 아닌 mixed trade-off")
    metrics = [
        ("Reclaim proxy", "30–32% ↓", TEAL),
        ("GC proxy", "22–23% ↓", CYAN),
        ("Peak RSS", "3–11% ↓", ORANGE),
        ("Heap stability Sₕ", "1–8% ↓", RED),
    ]
    for i, (name, value, col) in enumerate(metrics):
        x = 0.7 + i * 3.12
        add_rect(s, x, 1.55, 2.75, 2.0, NAVY2, col, True)
        add_text(s, name, x + 0.2, 1.88, 2.35, 0.4, 15, MUTED, True, PP_ALIGN.CENTER)
        add_text(s, value, x + 0.2, 2.48, 2.35, 0.55, 27, col, True, PP_ALIGN.CENTER)
    add_bullets(s, [
        "bounded intervention이 footprint variation을 추가할 수 있음",
        "peak와 variability를 함께 보고해야 함",
        "단일 seed의 synthetic output은 device performance 근거가 아님",
        "현재 결과의 역할: controller sanity check + 반증 가능한 trade-off 제시",
    ], 1.0, 4.2, 11.3, 2.0, 18)

    # 10. Reviewer fixes
    s = prs.slides.add_slide(blank); add_bg(s, 10)
    add_title(s, "리뷰 피드백 반영: 무엇이 바뀌었나")
    fixes = [
        ("신뢰성", "실측/배포 주장 제거\nsimulation scope 명시", TEAL),
        ("인용", "CXL·flash GC·PIM 오인용 제거\n공식 Android/Linux 자료로 교체", CYAN),
        ("수식", "Amax leakage 제거\nonline trailing p95 적용", ORANGE),
        ("Metric", "Sₕ 정의를 코드와 통일\npeak RSS 병기", TEAL),
        ("구조", "결과를 Evaluation으로 이동\ncase study → validation plan", CYAN),
        ("재현성", "raw schema·baseline·통계·장기시험 요구사항 명시", ORANGE),
    ]
    for i, (t, b, col) in enumerate(fixes):
        x = 0.7 + (i % 3) * 4.15
        y = 1.35 + (i // 3) * 2.45
        card(s, t, b, x, y, 3.8, 1.95, col)

    # 11. Validation roadmap
    s = prs.slides.add_slide(blank); add_bg(s, 11)
    add_title(s, "출판 가능한 device study로 가는 로드맵")
    phases = [
        ("A", "Implementation", "AOSP/ART/kernel patch\nSELinux + build provenance", TEAL),
        ("B", "Baselines", "LMKD/PSI · memory.reclaim\nDAMON · MGLRU · ablations", CYAN),
        ("C", "Devices", "4 memory tiers\n2+ SoC vendors", ORANGE),
        ("D", "Evidence", "raw Perfetto + per-run rows\nbootstrap/Wilcoxon + effect size", RED),
    ]
    for i, (n, t, b, col) in enumerate(phases):
        x = 0.7 + i * 3.12
        add_rect(s, x, 1.55, 2.75, 3.45, NAVY2, col, True)
        add_text(s, n, x + 0.2, 1.8, 0.4, 0.4, 20, col, True)
        add_text(s, t, x + 0.65, 1.78, 1.8, 0.4, 18, WHITE, True)
        add_text(s, b, x + 0.28, 2.55, 2.2, 1.25, 15, MUTED, False, PP_ALIGN.CENTER)
    add_text(s, "추가 필수: 6–24h stability · energy/thermal · refault · correctness · fairness",
             0.95, 5.65, 11.4, 0.45, 19, WHITE, True, PP_ALIGN.CENTER)

    # 12. Conclusion
    s = prs.slides.add_slide(blank); add_bg(s, 12)
    add_title(s, "결론")
    add_text(s, "HMS의 현재 기여", 0.8, 1.45, 4.2, 0.5, 24, TEAL, True)
    add_bullets(s, [
        "managed/native/page 관측을 잇는 state model",
        "online normalization + bounded/fair trigger",
        "seeded simulation과 검증 계획",
    ], 0.9, 2.15, 5.4, 2.25, 19)
    add_text(s, "HMS가 아직 증명하지 않은 것", 6.8, 1.45, 5.6, 0.5, 24, RED, True)
    add_bullets(s, [
        "상용 Android 구현 및 성능",
        "생산 overhead / compatibility / energy",
        "기존 proactive reclaim 대비 우월성",
    ], 6.9, 2.15, 5.4, 2.25, 19)
    add_rect(s, 0.8, 5.25, 11.7, 0.95, NAVY2, TEAL, True)
    add_text(s, "다음 단계: auditable implementation + raw evidence",
             1.1, 5.52, 11.1, 0.35, 23, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, "감사합니다  ·  Questions?", 4.2, 6.55, 4.9, 0.35, 18, MUTED, True, PP_ALIGN.CENTER)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
